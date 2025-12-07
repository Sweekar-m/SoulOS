from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import sys
import subprocess
from dotenv import load_dotenv
load_dotenv()  # loads .env automatically
os.environ["GOOGLE_API_KEY"] = os.getenv("GOOGLE_API_KEY")  
try:
    from google import genai
    print("✅ Gemini module imported successfully.")
except Exception as e:
    print(f"❌ Import error: {e}")
topic=""
slide=""
def ppt_file_fun(file_name, slide_num, Topic, save_path="Documents"):
    # Save path h
    global topic ,slide
    slide=slide_num
    topic=Topic
    if save_path.lower() == "documents":
        save_dir = os.path.join(os.path.expanduser("~"), "Documents")
    else:
        save_dir = save_path
    os.makedirs(save_dir, exist_ok=True)
    full_path = os.path.join(save_dir, file_name)

    # Gemini prompt
    def code_generate():
        global topic ,slide
        print (topic)
        prompt = f'''
Give me the information on the topic [topic] in the same format up to [slide] slide(s).
Each slide can have paragraph or point-wise information. If point-wise, include 4 points are  enough for one slide .
The last slide should contain the conclusion. Give the result in the exact format as shown below, without any explanations.

If the content is in paragraph form, enclose it in [" "] brackets. Do not include explanations or formatting like **bold**.
Output the result as a Python dictionary with the variable name `topic_details` and in the format given below.

Example:
topic_details = {{
    "title": "Chat GPT",
    "subtitle": "Understanding and Utilizing the Power of AI Conversations",
    "content_slides": [
        {{
            "title": "What is Chat GPT?",
            "content": [
                "* A cutting-edge language model developed by OpenAI.",
                "* Designed for engaging in natural and interactive conversations.",
                "* Leverages deep learning techniques for advanced text generation.",
                "* Capable of understanding context and providing relevant responses.",
                "* Continuously learns and adapts based on user interactions.",
                "* Can be used for various applications, from customer service to content creation."
            ]
        }}
    ]
}}
'''
        prompt=prompt.replace("[topic]",str(topic))
        prompt=prompt.replace("[slide]",str(slide))
        client = genai.Client()
        response = client.models.generate_content(
            model="gemini-2.5-flash", contents=prompt
        )
        return response.text

    # Formatting functions
    def set_title_slide(slide, title_text, subtitle_text):
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        title_shape.text = title_text
        subtitle_shape.text = subtitle_text
        title_tf = title_shape.text_frame
        subtitle_tf = subtitle_shape.text_frame
        title_tf.paragraphs[0].font.size = Pt(48)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = RGBColor(250, 250, 250)
        subtitle_tf.paragraphs[0].font.size = Pt(28)
        subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    def set_title(slide, title_text, font_size):
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        text_frame = title_shape.text_frame
        text_frame.text = title_text
        p = text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def set_body_content(slide, body_text_list, font_size):
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        for item in body_text_list:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT

    def set_background(slide, start_color, end_color):
        fill = slide.background.fill
        fill.gradient()
        stops = fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = start_color
        stops[1].position = 1.0
        stops[1].color.rgb = end_color

    # Presentation creator
    def create_presentation(details):
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
        set_title_slide(slide_1, details["title"], details["subtitle"] or "")
        set_background(slide_1, RGBColor(10, 10, 20), RGBColor(30, 30, 50))

        for slide_detail in details["content_slides"]:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            set_title(slide, slide_detail["title"], 45)
            set_body_content(slide, slide_detail["content"], 25)
            set_background(slide, RGBColor(10, 10, 20), RGBColor(30, 30, 50))

        prs.save(full_path)
        print(f"✅ Presentation saved to: {full_path}")
        try:
            os.startfile(full_path)
        except Exception as e:
            print(f"⚠️ Could not open PowerPoint automatically: {e}")

    # Run
    
    generated_info=code_generate()
    generated_info=generated_info.replace("```python", "")
    generated_info=generated_info.replace("```", "")
    try:
        local_scope = {}
        exec(generated_info, {}, local_scope)
        topic_details = local_scope["topic_details"]
        create_presentation(topic_details)
    except Exception as e:
        print("❌ Error extracting topic_details:", e)
        print("📜 Raw text from Gemini:\n", generated_info)

# Run the function
